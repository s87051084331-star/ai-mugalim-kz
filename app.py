from fastapi import FastAPI, UploadFile, File, HTTPException, Request
from fastapi.responses import FileResponse, Response
from fastapi.staticfiles import StaticFiles
from starlette.middleware.sessions import SessionMiddleware
from pydantic import BaseModel
from pathlib import Path
from datetime import datetime, timedelta
from io import BytesIO
import io, os, csv, re, json, secrets, base64, sqlite3, hashlib, hmac

app = FastAPI(title="AI Мұғалім көмекшісі — ZEREK Education")
app.add_middleware(SessionMiddleware, secret_key=os.getenv("SESSION_SECRET", secrets.token_hex(32)), same_site="lax", https_only=True)
BASE = Path(__file__).parent


# ---------------- PostgreSQL persistent access control ----------------
DATABASE_URL=os.getenv("DATABASE_URL","").strip()
if not DATABASE_URL:
    raise RuntimeError("DATABASE_URL is required. Add Render Postgres Internal Database URL to Environment Variables.")

import psycopg
from psycopg.rows import dict_row

def db():
    return psycopg.connect(DATABASE_URL, row_factory=dict_row)

def hash_password(password,salt=None):
    salt=salt or secrets.token_hex(16)
    dk=hashlib.pbkdf2_hmac("sha256",password.encode(),salt.encode(),200000)
    return salt+"$"+dk.hex()

def verify_password(password,stored):
    try:
        salt,digest=stored.split("$",1)
        return hmac.compare_digest(hash_password(password,salt).split("$",1)[1],digest)
    except Exception:return False

def init_db():
    with db() as c:
        with c.cursor() as cur:
            cur.execute("""CREATE TABLE IF NOT EXISTS users(
              id BIGSERIAL PRIMARY KEY,
              email TEXT UNIQUE NOT NULL,
              name TEXT NOT NULL,
              password_hash TEXT NOT NULL,
              status TEXT NOT NULL DEFAULT 'pending',
              role TEXT NOT NULL DEFAULT 'teacher',
              plan TEXT NOT NULL DEFAULT 'free',
              subscription_until TIMESTAMPTZ,
              openai_limit INTEGER NOT NULL DEFAULT 10,
              openai_used INTEGER NOT NULL DEFAULT 0,
              gemini_limit INTEGER NOT NULL DEFAULT 10,
              gemini_used INTEGER NOT NULL DEFAULT 0,
              usage_date DATE DEFAULT CURRENT_DATE,
              ai_provider TEXT NOT NULL DEFAULT 'auto',
              credits INTEGER NOT NULL DEFAULT 0,
              created_at TIMESTAMPTZ DEFAULT NOW()
            )""")
            cur.execute("""CREATE TABLE IF NOT EXISTS lesson_archive(
              id BIGSERIAL PRIMARY KEY,
              user_id BIGINT REFERENCES users(id) ON DELETE CASCADE,
              title TEXT NOT NULL,
              subject TEXT DEFAULT '',
              class_name TEXT DEFAULT '',
              topic TEXT DEFAULT '',
              learning_objective TEXT DEFAULT '',
              kmj_json JSONB NOT NULL DEFAULT '{}'::jsonb,
              created_at TIMESTAMPTZ DEFAULT NOW(),
              updated_at TIMESTAMPTZ DEFAULT NOW()
            )""")
            cur.execute("""CREATE TABLE IF NOT EXISTS kmj_tasks(
              id BIGSERIAL PRIMARY KEY,
              user_id BIGINT REFERENCES users(id) ON DELETE CASCADE,
              archive_id BIGINT REFERENCES lesson_archive(id) ON DELETE SET NULL,
              stage TEXT DEFAULT '',
              work_type TEXT DEFAULT 'individual',
              title TEXT NOT NULL,
              task_text TEXT NOT NULL,
              minutes INTEGER DEFAULT 5,
              points INTEGER DEFAULT 1,
              descriptor TEXT DEFAULT '',
              assessment TEXT DEFAULT '',
              resource_name TEXT DEFAULT '',
              resource_text TEXT DEFAULT '',
              created_at TIMESTAMPTZ DEFAULT NOW()
            )""")
            cur.execute("""CREATE TABLE IF NOT EXISTS payment_requests(
              id BIGSERIAL PRIMARY KEY,
              user_id BIGINT REFERENCES users(id) ON DELETE CASCADE,
              package_code TEXT NOT NULL,
              amount_kzt INTEGER NOT NULL,
              status TEXT NOT NULL DEFAULT 'pending',
              note TEXT DEFAULT '',
              created_at TIMESTAMPTZ DEFAULT NOW(),
              approved_at TIMESTAMPTZ
            )""")
            cur.execute("""CREATE TABLE IF NOT EXISTS teacher_resources(
              id BIGSERIAL PRIMARY KEY,
              user_id BIGINT REFERENCES users(id) ON DELETE CASCADE,
              name TEXT NOT NULL,
              resource_type TEXT DEFAULT 'link',
              url TEXT DEFAULT '',
              note TEXT DEFAULT '',
              created_at TIMESTAMPTZ DEFAULT NOW()
            )""")
            cur.execute("""CREATE TABLE IF NOT EXISTS zerek_task_bank(
              id BIGSERIAL PRIMARY KEY,
              user_id BIGINT REFERENCES users(id) ON DELETE SET NULL,
              subject TEXT NOT NULL, grade TEXT DEFAULT '', learning_objective TEXT DEFAULT '',
              title TEXT NOT NULL, task_text TEXT NOT NULL, level TEXT DEFAULT 'A',
              points INTEGER DEFAULT 1, descriptor TEXT DEFAULT '', answer_key TEXT DEFAULT '',
              is_public BOOLEAN DEFAULT TRUE, created_at TIMESTAMPTZ DEFAULT NOW()
            )""")
            cur.execute("""CREATE TABLE IF NOT EXISTS community_kmj(
              id BIGSERIAL PRIMARY KEY,
              user_id BIGINT REFERENCES users(id) ON DELETE SET NULL,
              title TEXT NOT NULL, subject TEXT DEFAULT '', class_name TEXT DEFAULT '',
              learning_objective TEXT DEFAULT '', kmj_json JSONB NOT NULL,
              is_public BOOLEAN DEFAULT TRUE, created_at TIMESTAMPTZ DEFAULT NOW()
            )""")
            cur.execute("""CREATE TABLE IF NOT EXISTS teacher_classes(
              id BIGSERIAL PRIMARY KEY,
              user_id BIGINT REFERENCES users(id) ON DELETE CASCADE,
              name TEXT NOT NULL, school_year TEXT DEFAULT '',
              created_at TIMESTAMPTZ DEFAULT NOW(),
              UNIQUE(user_id,name,school_year)
            )""")
            cur.execute("""CREATE TABLE IF NOT EXISTS class_students(
              id BIGSERIAL PRIMARY KEY,
              class_id BIGINT REFERENCES teacher_classes(id) ON DELETE CASCADE,
              name TEXT NOT NULL,
              student_code TEXT DEFAULT '',
              created_at TIMESTAMPTZ DEFAULT NOW(),
              UNIQUE(class_id,name)
            )""")
            cur.execute("""CREATE TABLE IF NOT EXISTS ai_usage(
              id BIGSERIAL PRIMARY KEY,
              user_id BIGINT REFERENCES users(id) ON DELETE CASCADE,
              provider TEXT NOT NULL,
              feature TEXT NOT NULL DEFAULT 'general',
              units INTEGER NOT NULL DEFAULT 1,
              created_at TIMESTAMPTZ DEFAULT NOW()
            )""")
            cur.execute("""CREATE TABLE IF NOT EXISTS plans(
              code TEXT PRIMARY KEY,
              name TEXT NOT NULL,
              openai_limit INTEGER NOT NULL,
              gemini_limit INTEGER NOT NULL,
              credits INTEGER NOT NULL DEFAULT 0
            )""")
            cur.executemany("""INSERT INTO plans(code,name,openai_limit,gemini_limit,credits)
              VALUES(%s,%s,%s,%s,%s) ON CONFLICT(code) DO UPDATE SET
              name=EXCLUDED.name,openai_limit=EXCLUDED.openai_limit,gemini_limit=EXCLUDED.gemini_limit,credits=EXCLUDED.credits""",[
              ("free","Free",10,10,0),("standard","Standard",200,100,200),("pro","Pro",1000,500,1200)
            ])
            ae=os.getenv("ADMIN_EMAIL","").strip().lower();ap=os.getenv("ADMIN_PASSWORD","").strip()
            if ae and ap:
                cur.execute("SELECT id FROM users WHERE email=%s",(ae,));u=cur.fetchone()
                if u:cur.execute("UPDATE users SET role='admin',status='approved' WHERE email=%s",(ae,))
                else:cur.execute("INSERT INTO users(email,name,password_hash,status,role,plan,openai_limit,gemini_limit) VALUES(%s,%s,%s,'approved','admin','pro',10000,10000)",(ae,"Administrator",hash_password(ap)))
init_db()

def reset_usage(u):
    if str(u.get("usage_date")) != datetime.now().strftime("%Y-%m-%d"):
        with db() as c:
            with c.cursor() as cur:cur.execute("UPDATE users SET openai_used=0,gemini_used=0,usage_date=CURRENT_DATE WHERE id=%s",(u["id"],))
        u=dict(u);u["openai_used"]=0;u["gemini_used"]=0;u["usage_date"]=datetime.now().date()
    return dict(u)

def auth_user(request:Request):
    uid=request.session.get("uid")
    if not uid:raise HTTPException(401,"Кіру қажет.")
    with db() as c:
        with c.cursor() as cur:cur.execute("SELECT * FROM users WHERE id=%s",(uid,));u=cur.fetchone()
    if not u or u["status"]!="approved":raise HTTPException(403,"Рұқсат жоқ.")
    return reset_usage(u)

class RegisterBody(BaseModel): email:str; name:str; password:str
class LoginBody(BaseModel): email:str; password:str
class ApprovalBody(BaseModel): user_id:int; approve:bool
class AiPrefBody(BaseModel): provider:str="auto"
class UserPlanBody(BaseModel): user_id:int; plan:str
class ProviderLimitBody(BaseModel): user_id:int; openai_limit:int; gemini_limit:int

@app.post("/api/auth/register")
def register(body:RegisterBody):
    email=body.email.strip().lower()
    if "@" not in email or len(body.password)<6 or not body.name.strip():raise HTTPException(400,"Email, аты-жөні және кемінде 6 таңбалы құпиясөз қажет.")
    try:
        with db() as c:
            with c.cursor() as cur:cur.execute("INSERT INTO users(email,name,password_hash) VALUES(%s,%s,%s)",(email,body.name.strip(),hash_password(body.password)))
    except psycopg.errors.UniqueViolation:raise HTTPException(409,"Бұл email бұрын тіркелген.")
    return {"ok":True,"status":"pending"}

@app.post("/api/auth/login")
def login(body:LoginBody,request:Request):
    with db() as c:
        with c.cursor() as cur:cur.execute("SELECT * FROM users WHERE email=%s",(body.email.strip().lower(),));u=cur.fetchone()
    if not u or not verify_password(body.password,u["password_hash"]):raise HTTPException(401,"Email немесе құпиясөз қате.")
    if u["status"]!="approved":raise HTTPException(403,"Аккаунт әлі әкімші тарапынан мақұлданбаған.")
    request.session["uid"]=u["id"];return {"ok":True,"user":{"id":u["id"],"email":u["email"],"name":u["name"],"role":u["role"]}}

@app.post("/api/auth/logout")
def logout(request:Request):request.session.clear();return {"ok":True}

@app.get("/api/auth/me")
def me(request:Request):
    try:u=auth_user(request)
    except HTTPException:return {"authenticated":False}
    return {"authenticated":True,"user":{"id":u["id"],"email":u["email"],"name":u["name"],"role":u["role"],"plan":u["plan"]}}

def require_admin(request):
    u=auth_user(request)
    if u["role"]!="admin":raise HTTPException(403,"Әкімші рұқсаты қажет.")
    return u

@app.get("/api/user/ai-settings")
def user_ai_settings(request:Request):
    u=auth_user(request)
    return {"provider":u["ai_provider"],"plan":u["plan"],"credits":u["credits"],
      "openai":{"used":u["openai_used"],"limit":u["openai_limit"]},
      "gemini":{"used":u["gemini_used"],"limit":u["gemini_limit"]}}

@app.post("/api/user/ai-settings")
def user_ai_save(body:AiPrefBody,request:Request):
    u=auth_user(request);p=body.provider if body.provider in ("auto","gemini","openai") else "auto"
    with db() as c:
        with c.cursor() as cur:cur.execute("UPDATE users SET ai_provider=%s WHERE id=%s",(p,u["id"]))
    return {"ok":True}

@app.get("/api/admin/users")
def admin_users(request:Request):
    require_admin(request)
    with db() as c:
        with c.cursor() as cur:cur.execute("""SELECT id,email,name,status,role,plan,subscription_until,openai_limit,openai_used,gemini_limit,gemini_used,credits,created_at FROM users ORDER BY id DESC""");rows=cur.fetchall()
    return {"users":rows}

@app.post("/api/admin/approve")
def admin_approve(body:ApprovalBody,request:Request):
    require_admin(request)
    with db() as c:
        with c.cursor() as cur:cur.execute("UPDATE users SET status=%s WHERE id=%s",("approved" if body.approve else "rejected",body.user_id))
    return {"ok":True}

@app.post("/api/admin/user-plan")
def admin_plan(body:UserPlanBody,request:Request):
    require_admin(request);code=body.plan if body.plan in ("free","standard","pro") else "free"
    with db() as c:
        with c.cursor() as cur:
            cur.execute("SELECT * FROM plans WHERE code=%s",(code,));p=cur.fetchone()
            cur.execute("""UPDATE users SET plan=%s,openai_limit=%s,gemini_limit=%s,credits=%s,
              openai_used=0,gemini_used=0,usage_date=CURRENT_DATE WHERE id=%s""",
              (code,p["openai_limit"],p["gemini_limit"],p["credits"],body.user_id))
    return {"ok":True}

@app.post("/api/admin/provider-limits")
def admin_provider_limits(body:ProviderLimitBody,request:Request):
    require_admin(request)
    with db() as c:
        with c.cursor() as cur:cur.execute("UPDATE users SET openai_limit=%s,gemini_limit=%s WHERE id=%s",(max(0,body.openai_limit),max(0,body.gemini_limit),body.user_id))
    return {"ok":True}

def charge_provider(user_id,provider,feature="general"):
    col="openai_used" if provider=="openai" else "gemini_used"
    with db() as c:
        with c.cursor() as cur:
            cur.execute(f"UPDATE users SET {col}={col}+1 WHERE id=%s",(user_id,))
            cur.execute("INSERT INTO ai_usage(user_id,provider,feature,units) VALUES(%s,%s,%s,1)",(user_id,provider,feature))

@app.get("/api/admin/usage")
def admin_usage(request:Request):
    require_admin(request)
    with db() as c:
        with c.cursor() as cur:
            cur.execute("""SELECT u.email,a.provider,a.feature,COUNT(*) AS uses,MAX(a.created_at) AS last_used
              FROM ai_usage a JOIN users u ON u.id=a.user_id GROUP BY u.email,a.provider,a.feature ORDER BY last_used DESC LIMIT 500""")
            rows=cur.fetchall()
    return {"usage":rows}


class KmjGenerateBody(BaseModel):
    subject:str=""
    class_name:str=""
    section:str=""
    topic:str=""
    learning_objective:str=""
    lesson_goal:str=""
    duration:int=45
    value:str=""
    weekly_quote:str=""
    resources:list=[]
    extra_instruction:str=""
    provider:str="auto"

class ArchiveSaveBody(BaseModel):
    title:str
    subject:str=""
    class_name:str=""
    topic:str=""
    learning_objective:str=""
    kmj:dict

class ResourceBody(BaseModel):
    name:str
    resource_type:str="link"
    url:str=""
    note:str=""

@app.post("/api/kmj/generate130")
async def generate_kmj130(body:KmjGenerateBody,request:Request):
    u=auth_user(request)
    prompt=f"""Қазақстан мұғаліміне ҚМЖ жобасын жаса.
Кесте құрылымы міндетті түрде 5 бағаннан тұрады:
1) Сабақтың кезеңі / Уақыт
2) Педагогтің іс-әрекеті
3) Оқушының іс-әрекеті
4) Бағалау
5) Ресурстар.
Жоғарғы мәліметтер: пән, сынып, бөлім, сабақ тақырыбы, оқу мақсаты, сабақ мақсаты.
Құндылық және Аптаның дәйексөзі — мұғалім енгізген мәтінді өзгеріссіз сақта.
Мұғалім берген ресурстарды орынды кезеңдерге орналастыр, бірақ жаңа URL ойдан шығарма.
ПӘНГЕ БЕЙІМДЕУ МІНДЕТТІ:
- Математика/Алгебра/Геометрияда "есеп шығарады" деп жалпылама жазба: есептің нақты шартын, өрнегін немесе теңдеуін толық бер.
- Физикада нақты есеп, формула, өлшем бірліктерін; химияда теңдеу/есеп/тәжірибені; тіл пәндерінде оқылым/жазылым/айтылым тапсырмаларын нақты бер.
- Сабақ ортасында кемінде A/B/C деңгейлерін қамтитын нақты тапсырмалар болсын.
- Әр тапсырманың дескрипторы мен баллы бағалау бағанында анық көрсетілсін.
- Мұғалімге арналған жауап кілтін resources немесе teacher_action ішінде "Жауап кілті:" деп белгіле.
- Сабақ соңында нақты бекіту тапсырмасы, рефлексия және нақты үй тапсырмасы міндетті болсын.
- Уақыттардың қосындысы берілген сабақ ұзақтығына сәйкес келсін.
Тек JSON қайтар:
{{"subject":"","class_name":"","section":"","topic":"","learning_objective":"","lesson_goal":"",
"value":"","weekly_quote":"","duration":45,
"rows":[{{"stage":"","minutes":0,"teacher_action":"","student_action":"","assessment":"","resources":""}}]}}
Дерек:
Пән: {body.subject}
Сынып: {body.class_name}
Бөлім: {body.section}
Тақырып: {body.topic}
Оқу мақсаты: {body.learning_objective}
Сабақ мақсаты: {body.lesson_goal}
Ұзақтығы: {body.duration}
Құндылық: {body.value}
Аптаның дәйексөзі: {body.weekly_quote}
Ресурстар: {json.dumps(body.resources,ensure_ascii=False)}
Қосымша нұсқау: {body.extra_instruction}
"""
    ai=await ai_json(body.provider,prompt,u,"kmj_generate130")
    data=ai["data"]
    return {"ok":True,"mode":ai["provider_used"],"fallback":ai["fallback"],"data":data}

@app.post("/api/archive/save")
def archive_save(body:ArchiveSaveBody,request:Request):
    u=auth_user(request)
    with db() as c:
        with c.cursor() as cur:
            cur.execute("""INSERT INTO lesson_archive(user_id,title,subject,class_name,topic,learning_objective,kmj_json)
              VALUES(%s,%s,%s,%s,%s,%s,%s) RETURNING id""",
              (u["id"],body.title,body.subject,body.class_name,body.topic,body.learning_objective,json.dumps(body.kmj,ensure_ascii=False)))
            rid=cur.fetchone()["id"]
    return {"ok":True,"id":rid}

@app.get("/api/archive")
def archive_list(request:Request):
    u=auth_user(request)
    with db() as c:
        with c.cursor() as cur:
            cur.execute("""SELECT id,title,subject,class_name,topic,learning_objective,created_at,updated_at
              FROM lesson_archive WHERE user_id=%s ORDER BY updated_at DESC""",(u["id"],));rows=cur.fetchall()
    return {"items":rows}

@app.get("/api/archive/{item_id}")
def archive_get(item_id:int,request:Request):
    u=auth_user(request)
    with db() as c:
        with c.cursor() as cur:
            cur.execute("SELECT * FROM lesson_archive WHERE id=%s AND user_id=%s",(item_id,u["id"]));row=cur.fetchone()
    if not row:raise HTTPException(404,"Архив материалы табылмады.")
    return {"item":row}

@app.delete("/api/archive/{item_id}")
def archive_delete(item_id:int,request:Request):
    u=auth_user(request)
    with db() as c:
        with c.cursor() as cur:cur.execute("DELETE FROM lesson_archive WHERE id=%s AND user_id=%s",(item_id,u["id"]))
    return {"ok":True}

@app.post("/api/resources")
def resource_add(body:ResourceBody,request:Request):
    u=auth_user(request)
    with db() as c:
        with c.cursor() as cur:
            cur.execute("""INSERT INTO teacher_resources(user_id,name,resource_type,url,note)
              VALUES(%s,%s,%s,%s,%s) RETURNING id""",(u["id"],body.name,body.resource_type,body.url,body.note));rid=cur.fetchone()["id"]
    return {"ok":True,"id":rid}

@app.get("/api/resources")
def resources_list(request:Request):
    u=auth_user(request)
    with db() as c:
        with c.cursor() as cur:cur.execute("SELECT * FROM teacher_resources WHERE user_id=%s ORDER BY id DESC",(u["id"],));rows=cur.fetchall()
    return {"items":rows}



class KmjTaskBody(BaseModel):
    stage:str=""
    work_type:str="individual"
    title:str
    task_text:str
    minutes:int=5
    points:int=1
    descriptor:str=""
    assessment:str=""
    resource_name:str=""
    resource_text:str=""

class PaymentRequestBody(BaseModel):
    package_code:str
    note:str=""

class PaymentApproveBody(BaseModel):
    payment_id:int
    approve:bool

PACKAGE_CATALOG={
 "standard":{"name":"Standard","amount_kzt":2990,"openai_add":200,"gemini_add":100,"credits_add":200},
 "pro":{"name":"Pro","amount_kzt":6990,"openai_add":1000,"gemini_add":500,"credits_add":1200},
 "openai100":{"name":"OpenAI +100","amount_kzt":1490,"openai_add":100,"gemini_add":0,"credits_add":0},
 "gemini100":{"name":"Gemini +100","amount_kzt":990,"openai_add":0,"gemini_add":100,"credits_add":0}
}

@app.post("/api/kmj/tasks")
def kmj_task_add(body:KmjTaskBody,request:Request):
    u=auth_user(request)
    with db() as c:
        with c.cursor() as cur:
            cur.execute("""INSERT INTO kmj_tasks(user_id,stage,work_type,title,task_text,minutes,points,descriptor,assessment,resource_name,resource_text)
              VALUES(%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s) RETURNING id""",
              (u["id"],body.stage,body.work_type,body.title,body.task_text,max(1,body.minutes),max(0,body.points),body.descriptor,body.assessment,body.resource_name,body.resource_text))
            rid=cur.fetchone()["id"]
    return {"ok":True,"id":rid}

@app.get("/api/kmj/tasks")
def kmj_tasks(request:Request):
    u=auth_user(request)
    with db() as c:
        with c.cursor() as cur:cur.execute("SELECT * FROM kmj_tasks WHERE user_id=%s ORDER BY id DESC",(u["id"],));rows=cur.fetchall()
    return {"items":rows}

@app.delete("/api/kmj/tasks/{task_id}")
def kmj_task_delete(task_id:int,request:Request):
    u=auth_user(request)
    with db() as c:
        with c.cursor() as cur:cur.execute("DELETE FROM kmj_tasks WHERE id=%s AND user_id=%s",(task_id,u["id"]))
    return {"ok":True}

@app.get("/api/packages")
def packages(request:Request):
    auth_user(request)
    return {"packages":[{"code":k,**v} for k,v in PACKAGE_CATALOG.items()]}

@app.post("/api/payments/request")
def payment_request(body:PaymentRequestBody,request:Request):
    u=auth_user(request);p=PACKAGE_CATALOG.get(body.package_code)
    if not p:raise HTTPException(400,"Пакет табылмады.")
    with db() as c:
        with c.cursor() as cur:
            cur.execute("""INSERT INTO payment_requests(user_id,package_code,amount_kzt,note)
              VALUES(%s,%s,%s,%s) RETURNING id""",(u["id"],body.package_code,p["amount_kzt"],body.note));pid=cur.fetchone()["id"]
    return {"ok":True,"payment_id":pid,"status":"pending","amount_kzt":p["amount_kzt"]}

@app.get("/api/payments/mine")
def my_payments(request:Request):
    u=auth_user(request)
    with db() as c:
        with c.cursor() as cur:cur.execute("SELECT * FROM payment_requests WHERE user_id=%s ORDER BY id DESC",(u["id"],));rows=cur.fetchall()
    return {"items":rows}

@app.get("/api/admin/payments")
def admin_payments(request:Request):
    require_admin(request)
    with db() as c:
        with c.cursor() as cur:
            cur.execute("""SELECT p.*,u.email,u.name FROM payment_requests p JOIN users u ON u.id=p.user_id ORDER BY p.id DESC""");rows=cur.fetchall()
    return {"items":rows}

@app.post("/api/admin/payments/approve")
def admin_payment_approve(body:PaymentApproveBody,request:Request):
    require_admin(request)
    with db() as c:
        with c.cursor() as cur:
            cur.execute("SELECT * FROM payment_requests WHERE id=%s FOR UPDATE",(body.payment_id,));pay=cur.fetchone()
            if not pay:raise HTTPException(404,"Төлем сұранысы табылмады.")
            if pay["status"]!="pending":return {"ok":True,"status":pay["status"]}
            if not body.approve:
                cur.execute("UPDATE payment_requests SET status='rejected',approved_at=NOW() WHERE id=%s",(body.payment_id,))
                return {"ok":True,"status":"rejected"}
            p=PACKAGE_CATALOG.get(pay["package_code"])
            if not p:raise HTTPException(400,"Пакет конфигурациясы жоқ.")
            cur.execute("""UPDATE users SET openai_limit=openai_limit+%s,gemini_limit=gemini_limit+%s,credits=credits+%s,
              plan=CASE WHEN %s IN ('standard','pro') THEN %s ELSE plan END WHERE id=%s""",
              (p["openai_add"],p["gemini_add"],p["credits_add"],pay["package_code"],pay["package_code"],pay["user_id"]))
            cur.execute("UPDATE payment_requests SET status='approved',approved_at=NOW() WHERE id=%s",(body.payment_id,))
    return {"ok":True,"status":"approved"}




class ZerekTaskBankBody(BaseModel):
    subject:str; grade:str=""; learning_objective:str=""; title:str; task_text:str
    level:str="A"; points:int=1; descriptor:str=""; answer_key:str=""

class ZerekClassBody(BaseModel):
    name:str; school_year:str=""

@app.get("/api/zerek/task-bank")
def zerek_task_bank(request:Request,subject:str="",grade:str="",q:str=""):
    auth_user(request)
    sql="""SELECT b.*,u.name AS author_name FROM zerek_task_bank b LEFT JOIN users u ON u.id=b.user_id WHERE b.is_public=TRUE"""
    vals=[]
    if subject: sql+=" AND LOWER(b.subject)=LOWER(%s)";vals.append(subject)
    if grade: sql+=" AND b.grade=%s";vals.append(grade)
    if q: sql+=" AND (b.title ILIKE %s OR b.task_text ILIKE %s OR b.learning_objective ILIKE %s)";vals += ["%"+q+"%"]*3
    sql+=" ORDER BY b.id DESC LIMIT 100"
    with db() as c:
        with c.cursor() as cur:cur.execute(sql,vals);rows=cur.fetchall()
    return {"items":rows}

@app.post("/api/zerek/task-bank")
def zerek_task_bank_add(body:ZerekTaskBankBody,request:Request):
    u=auth_user(request)
    if not body.title.strip() or not body.task_text.strip():raise HTTPException(400,"Тапсырма атауы мен мәтіні қажет.")
    with db() as c:
        with c.cursor() as cur:
            cur.execute("""INSERT INTO zerek_task_bank(user_id,subject,grade,learning_objective,title,task_text,level,points,descriptor,answer_key)
            VALUES(%s,%s,%s,%s,%s,%s,%s,%s,%s,%s) RETURNING id""",(u["id"],body.subject,body.grade,body.learning_objective,body.title,body.task_text,body.level,max(1,body.points),body.descriptor,body.answer_key))
            rid=cur.fetchone()["id"]
    return {"ok":True,"id":rid}

@app.get("/api/zerek/community-kmj")
def zerek_community_kmj(request:Request,subject:str="",grade:str="",q:str=""):
    auth_user(request)
    sql="""SELECT k.id,k.title,k.subject,k.class_name,k.learning_objective,k.created_at,u.name AS author_name
    FROM community_kmj k LEFT JOIN users u ON u.id=k.user_id WHERE k.is_public=TRUE"""
    vals=[]
    if subject:sql+=" AND LOWER(k.subject)=LOWER(%s)";vals.append(subject)
    if grade:sql+=" AND k.class_name=%s";vals.append(grade)
    if q:sql+=" AND (k.title ILIKE %s OR k.learning_objective ILIKE %s)";vals += ["%"+q+"%"]*2
    sql+=" ORDER BY k.id DESC LIMIT 100"
    with db() as c:
        with c.cursor() as cur:cur.execute(sql,vals);rows=cur.fetchall()
    return {"items":rows}

@app.post("/api/zerek/community-kmj/{item_id}/copy")
def zerek_community_copy(item_id:int,request:Request):
    u=auth_user(request)
    with db() as c:
        with c.cursor() as cur:
            cur.execute("SELECT * FROM community_kmj WHERE id=%s AND is_public=TRUE",(item_id,));x=cur.fetchone()
            if not x:raise HTTPException(404,"ҚМЖ табылмады.")
            cur.execute("""INSERT INTO lesson_archive(user_id,title,subject,class_name,learning_objective,kmj_json)
            VALUES(%s,%s,%s,%s,%s,%s) RETURNING id""",(u["id"],x["title"],x["subject"],x["class_name"],x["learning_objective"],json.dumps(x["kmj_json"],ensure_ascii=False) if not isinstance(x["kmj_json"],str) else x["kmj_json"]))
            rid=cur.fetchone()["id"]
    return {"ok":True,"id":rid}

@app.get("/api/zerek/classes")
def zerek_classes(request:Request):
    u=auth_user(request)
    with db() as c:
        with c.cursor() as cur:
            cur.execute("""SELECT c.*,(SELECT COUNT(*) FROM class_students s WHERE s.class_id=c.id) AS student_count FROM teacher_classes c WHERE c.user_id=%s ORDER BY c.name""",(u["id"],));rows=cur.fetchall()
    return {"items":rows}

@app.post("/api/zerek/classes")
def zerek_class_add(body:ZerekClassBody,request:Request):
    u=auth_user(request)
    if not body.name.strip():raise HTTPException(400,"Сынып атауын енгізіңіз.")
    with db() as c:
        with c.cursor() as cur:
            cur.execute("""INSERT INTO teacher_classes(user_id,name,school_year) VALUES(%s,%s,%s)
            ON CONFLICT(user_id,name,school_year) DO UPDATE SET name=EXCLUDED.name RETURNING id""",(u["id"],body.name.strip(),body.school_year.strip()))
            rid=cur.fetchone()["id"]
    return {"ok":True,"id":rid}




class ZerekStudentBody(BaseModel):
    name:str

def zerek_owned_class(cur,class_id,user_id):
    cur.execute("SELECT id,name,school_year FROM teacher_classes WHERE id=%s AND user_id=%s",(class_id,user_id))
    row=cur.fetchone()
    if not row: raise HTTPException(404,"Сынып табылмады.")
    return row

@app.get("/api/zerek/classes/{class_id}/students")
def zerek_students(class_id:int,request:Request):
    u=auth_user(request)
    with db() as c:
        with c.cursor() as cur:
            zerek_owned_class(cur,class_id,u["id"])
            cur.execute("SELECT * FROM class_students WHERE class_id=%s ORDER BY name",(class_id,))
            rows=cur.fetchall()
    return {"items":rows}

@app.post("/api/zerek/classes/{class_id}/students")
def zerek_student_add(class_id:int,body:ZerekStudentBody,request:Request):
    u=auth_user(request);name=body.name.strip()
    if not name: raise HTTPException(400,"Оқушы аты-жөні қажет.")
    with db() as c:
        with c.cursor() as cur:
            zerek_owned_class(cur,class_id,u["id"])
            cur.execute("""INSERT INTO class_students(class_id,name) VALUES(%s,%s)
            ON CONFLICT(class_id,name) DO UPDATE SET name=EXCLUDED.name RETURNING id""",(class_id,name))
            sid=cur.fetchone()["id"]
            cur.execute("UPDATE class_students SET student_code=%s WHERE id=%s",("ST"+str(sid).zfill(4),sid))
    return {"ok":True,"id":sid}

@app.delete("/api/zerek/classes/{class_id}/students/{student_id}")
def zerek_student_delete(class_id:int,student_id:int,request:Request):
    u=auth_user(request)
    with db() as c:
        with c.cursor() as cur:
            zerek_owned_class(cur,class_id,u["id"])
            cur.execute("DELETE FROM class_students WHERE id=%s AND class_id=%s",(student_id,class_id))
    return {"ok":True}

@app.post("/api/zerek/classes/{class_id}/import")
async def zerek_students_import(class_id:int,request:Request,file:UploadFile=File(...)):
    u=auth_user(request);data=await file.read();names=parse_class(file.filename,data)
    with db() as c:
        with c.cursor() as cur:
            zerek_owned_class(cur,class_id,u["id"])
            for name in names:
                cur.execute("INSERT INTO class_students(class_id,name) VALUES(%s,%s) ON CONFLICT(class_id,name) DO NOTHING",(class_id,name))
            cur.execute("SELECT id FROM class_students WHERE class_id=%s AND COALESCE(student_code,'')=''",(class_id,))
            for row in cur.fetchall():
                cur.execute("UPDATE class_students SET student_code=%s WHERE id=%s",("ST"+str(row["id"]).zfill(4),row["id"]))
    return {"ok":True,"count":len(names)}


@app.delete("/api/zerek/classes/{class_id}")
def zerek_class_delete(class_id:int,request:Request):
    u=auth_user(request)
    with db() as c:
        with c.cursor() as cur:
            cur.execute("SELECT id,name FROM teacher_classes WHERE id=%s AND user_id=%s",(class_id,u["id"]))
            row=cur.fetchone()
            if not row: raise HTTPException(404,"Сынып табылмады.")
            cur.execute("DELETE FROM teacher_classes WHERE id=%s AND user_id=%s",(class_id,u["id"]))
    return {"ok":True,"name":row["name"]}

@app.post("/api/zerek/community-kmj")
def zerek_community_publish(body:ArchiveSaveBody,request:Request):
    u=auth_user(request)
    with db() as c:
        with c.cursor() as cur:
            cur.execute("""INSERT INTO community_kmj(user_id,title,subject,class_name,learning_objective,kmj_json)
            VALUES(%s,%s,%s,%s,%s,%s) RETURNING id""",(u["id"],body.title,body.subject,body.class_name,body.learning_objective,json.dumps(body.kmj,ensure_ascii=False)))
            rid=cur.fetchone()["id"]
    return {"ok":True,"id":rid}

@app.post("/api/kmj/export-docx")
def export_kmj_docx(body:ArchiveSaveBody,request:Request):
    auth_user(request)
    from docx import Document
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    import tempfile
    k=body.kmj
    d=Document()
    h=d.add_paragraph();h.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r=h.add_run("Қысқа мерзімді жоспар");r.bold=True;r.font.size=Pt(16)
    for label,key in [("Пән","subject"),("Сынып","class_name"),("Бөлім","section"),("Сабақ тақырыбы","topic"),("Оқу мақсаты","learning_objective"),("Сабақ мақсаты","lesson_goal"),("Құндылық","value"),("Аптаның дәйексөзі","weekly_quote")]:
        p=d.add_paragraph();p.add_run(label+": ").bold=True;p.add_run(str(k.get(key,"")))
    table=d.add_table(rows=1,cols=5);table.style="Table Grid"
    heads=["Сабақтың кезеңі / Уақыт","Педагогтің іс-әрекеті","Оқушының іс-әрекеті","Бағалау","Ресурстар"]
    for i,x in enumerate(heads):table.rows[0].cells[i].text=x
    for row in k.get("rows",[]):
        c=table.add_row().cells
        c[0].text=f'{row.get("stage","")} / {row.get("minutes","")} мин'
        c[1].text=str(row.get("teacher_action",""));c[2].text=str(row.get("student_action",""));c[3].text=str(row.get("assessment",""));c[4].text=str(row.get("resources",""))
    path=BASE/f"kmj_{secrets.token_hex(5)}.docx";d.save(path)
    return FileResponse(path,filename="ZEREK_KMJ.docx",media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document")


# ---------------- File parsing ----------------
def parse_file(name: str, data: bytes) -> str:
    ext = Path(name).suffix.lower()
    try:
        if ext in {".txt", ".csv"}:
            for enc in ("utf-8-sig","utf-8","cp1251"):
                try: return data.decode(enc)
                except Exception: pass
            return data.decode("utf-8", errors="ignore")
        if ext == ".docx":
            from docx import Document
            d = Document(io.BytesIO(data))
            parts = [p.text for p in d.paragraphs if p.text.strip()]
            for t in d.tables:
                for row in t.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
            return "\n".join(parts)
        if ext == ".pdf":
            from pypdf import PdfReader
            return "\n".join((p.extract_text() or "") for p in PdfReader(io.BytesIO(data)).pages)
        if ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
            parts=[]
            for ws in wb.worksheets:
                parts.append(f"[{ws.title}]")
                for row in ws.iter_rows(values_only=True):
                    parts.append(" | ".join("" if v is None else str(v) for v in row))
            return "\n".join(parts)
        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            return "\n".join(sh.text for sl in prs.slides for sh in sl.shapes if hasattr(sh, "text"))
        if ext in {".jpg",".jpeg",".png",".webp"}:
            return "Сурет қабылданды. Бұл нұсқада сурет OCR-ы қосылмаған."
        return "Файл форматы оқуға қолдау таппады."
    except Exception as e:
        return f"Файлды оқу қатесі: {e}"

@app.post("/api/parse")
async def parse(file: UploadFile = File(...)):
    data = await file.read()
    return {"message": f"{file.filename} өңделді", "text": parse_file(file.filename, data)}


# ---------------- KMJ image extraction ----------------
def _data_uri(blob: bytes, mime: str):
    return "data:"+mime+";base64,"+base64.b64encode(blob).decode("ascii")

def extract_images(name: str, data: bytes):
    ext=Path(name).suffix.lower()
    out=[]
    try:
        if ext==".docx":
            import zipfile as _zf
            with _zf.ZipFile(io.BytesIO(data)) as z:
                for n in z.namelist():
                    if n.startswith("word/media/"):
                        b=z.read(n); suf=Path(n).suffix.lower()
                        mime={".png":"image/png",".jpg":"image/jpeg",".jpeg":"image/jpeg",".gif":"image/gif",".webp":"image/webp"}.get(suf,"application/octet-stream")
                        out.append({"name":Path(n).name,"mime":mime,"data_url":_data_uri(b,mime),"source":"DOCX"})
        elif ext==".pptx":
            from pptx import Presentation
            prs=Presentation(io.BytesIO(data))
            seen=set()
            for si,slide in enumerate(prs.slides,1):
                for shape in slide.shapes:
                    if hasattr(shape,"image"):
                        b=shape.image.blob
                        key=hash(b)
                        if key in seen: continue
                        seen.add(key)
                        mime=shape.image.content_type or "image/png"
                        out.append({"name":f"slide-{si}-{shape.image.filename}","mime":mime,"data_url":_data_uri(b,mime),"source":f"PPTX {si}-слайд"})
        elif ext==".xlsx":
            from openpyxl import load_workbook
            wb=load_workbook(io.BytesIO(data))
            for ws in wb.worksheets:
                for i,img in enumerate(getattr(ws,"_images",[]),1):
                    try:
                        b=img._data()
                        fmt=(getattr(img,"format","png") or "png").lower()
                        mime="image/jpeg" if fmt in ("jpg","jpeg") else "image/png"
                        out.append({"name":f"{ws.title}-{i}.{fmt}","mime":mime,"data_url":_data_uri(b,mime),"source":f"XLSX {ws.title}"})
                    except Exception: pass
        elif ext==".pdf":
            # Extract embedded raster images where pypdf exposes them.
            from pypdf import PdfReader
            reader=PdfReader(io.BytesIO(data))
            for pi,page in enumerate(reader.pages,1):
                try:
                    for ii,img in enumerate(page.images,1):
                        b=img.data
                        nm=getattr(img,"name",f"page-{pi}-{ii}.png")
                        suf=Path(nm).suffix.lower()
                        mime={".jpg":"image/jpeg",".jpeg":"image/jpeg",".png":"image/png",".jp2":"image/jp2"}.get(suf,"image/png")
                        out.append({"name":nm,"mime":mime,"data_url":_data_uri(b,mime),"source":f"PDF {pi}-бет"})
                except Exception: pass
    except Exception:
        pass
    return out[:40]

@app.post("/api/kmj/images")
async def kmj_images(file: UploadFile=File(...)):
    data=await file.read()
    images=extract_images(file.filename,data)
    return {"count":len(images),"images":images}

# ---------------- Class import ----------------
def clean(v): return "" if v is None else str(v).strip()
def looks_name(s):
    s=clean(s)
    if len(s)<3 or s.isdigit(): return False
    return any(ch.isalpha() for ch in s) and s.lower() not in {"аты-жөні","оқушы","фио","name","student"}

def parse_class(name, data):
    ext=Path(name).suffix.lower(); rows=[]
    if ext==".xlsx":
        from openpyxl import load_workbook
        ws=load_workbook(io.BytesIO(data),data_only=True,read_only=True).active
        rows=[[clean(v) for v in row] for row in ws.iter_rows(values_only=True)]
    elif ext in {".csv",".txt"}:
        text=None
        for enc in ("utf-8-sig","utf-8","cp1251"):
            try: text=data.decode(enc); break
            except Exception: pass
        text=text or data.decode("utf-8",errors="ignore")
        try: dialect=csv.Sniffer().sniff(text[:4096],delimiters=";,|\t,")
        except Exception: dialect=csv.excel
        rows=[[clean(v) for v in r] for r in csv.reader(io.StringIO(text),dialect)]
    rows=[r for r in rows if any(r)]
    if not rows:return []
    header=[x.lower() for x in rows[0]]
    idx=None
    for i,h in enumerate(header):
        if any(k in h for k in ("аты-жөні","аты жөні","оқушы","фио","name","student")): idx=i;break
    start=1 if idx is not None else 0
    if idx is None:
        width=max(len(r) for r in rows)
        scores=[sum(looks_name(r[c] if c<len(r) else "") for r in rows[:100]) for c in range(width)]
        idx=max(range(width),key=lambda c:scores[c])
    names=[]
    for r in rows[start:]:
        if idx<len(r) and looks_name(r[idx]) and r[idx] not in names:names.append(r[idx])
    return names[:60]

@app.post("/api/import-class")
async def import_class(file: UploadFile = File(...)):
    data=await file.read(); names=parse_class(file.filename,data)
    return {"count":len(names),"students":names,"message":f"{len(names)} оқушы табылды"}

# ---------------- Selectable AI provider: Gemini / OpenAI ----------------
import base64
import asyncio

def _json_from_text(text: str):
    text=(text or "").strip()
    if text.startswith("```"):
        text=re.sub(r"^```(?:json)?\s*","",text)
        text=re.sub(r"\s*```$","",text)
    return json.loads(text)

async def gemini_text_json(prompt: str):
    key=os.getenv("GEMINI_API_KEY","").strip()
    if not key: raise HTTPException(400,"GEMINI_API_KEY Render Environment ішінде орнатылмаған.")
    import httpx
    model=os.getenv("GEMINI_MODEL","gemini-2.5-flash")
    url=f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={key}"
    payload={"contents":[{"parts":[{"text":prompt}]}],"generationConfig":{"responseMimeType":"application/json"}}
    async with httpx.AsyncClient(timeout=90) as client:
        r=await client.post(url,json=payload)
        if r.status_code>=400: raise HTTPException(r.status_code,f"Gemini: {r.text[:500]}")
        text=r.json()["candidates"][0]["content"]["parts"][0]["text"]
        return _json_from_text(text)

async def openai_text_json(prompt: str):
    key=os.getenv("OPENAI_API_KEY","").strip()
    if not key: raise HTTPException(400,"OPENAI_API_KEY Render Environment ішінде орнатылмаған.")
    import httpx
    model=os.getenv("OPENAI_MODEL","gpt-5.6")
    payload={"model":model,"input":prompt}
    async with httpx.AsyncClient(timeout=120) as client:
        r=await client.post("https://api.openai.com/v1/responses",headers={"Authorization":f"Bearer {key}","Content-Type":"application/json"},json=payload)
        if r.status_code>=400: raise HTTPException(r.status_code,f"OpenAI: {r.text[:500]}")
        data=r.json()
        chunks=[]
        for item in data.get("output",[]):
            if item.get("type")=="message":
                for c in item.get("content",[]):
                    if c.get("type")=="output_text": chunks.append(c.get("text",""))
        return _json_from_text("\n".join(chunks))

def _ai_err(provider,e):
    code=e.status_code if isinstance(e,HTTPException) else 503
    msg=str(e.detail if isinstance(e,HTTPException) else e).replace("\n"," ")[:700]
    return {"provider":provider,"status":code,"message":msg}

async def _retry_ai(provider: str, prompt: str, attempts: int=3):
    errs=[]
    for n in range(attempts):
        try:
            data=await (openai_text_json(prompt) if provider=="openai" else gemini_text_json(prompt))
            return {"ok":True,"data":data,"attempts":n+1,"errors":errs}
        except Exception as e:
            x=_ai_err(provider,e);x["attempt"]=n+1;errs.append(x)
            if x["status"] in (400,401,403,404): break
            if n<attempts-1: await asyncio.sleep(1.5*(2**n))
    return {"ok":False,"errors":errs}

async def ai_json(provider: str, prompt: str, user=None, feature="general"):
    requested=(provider or "auto").lower()
    if user:
        user=reset_usage(user)
        pref=(user.get("ai_provider") or "auto").lower()
        if requested=="auto":requested=pref
    if requested not in ("gemini","openai"):requested="openai"
    configured={"gemini":bool(os.getenv("GEMINI_API_KEY","").strip()),"openai":bool(os.getenv("OPENAI_API_KEY","").strip())}
    other="openai" if requested=="gemini" else "gemini";order=[requested]+([other] if configured[other] else [])
    errs=[]
    for p in order:
        if not configured[p]:continue
        if user:
            used=user[f"{p}_used"];limit=user[f"{p}_limit"]
            if used>=limit:
                errs.append({"provider":p,"status":429,"message":f"Жеке {p} лимиті аяқталды: {used}/{limit}"});continue
        r=await _retry_ai(p,prompt,3);errs+=r["errors"]
        if r["ok"]:
            if user:charge_provider(user["id"],p,feature)
            return {"provider_used":p,"requested_provider":requested,"fallback":p!=requested,"attempts":r["attempts"],"data":r["data"]}
    latest={}
    for e in errs:latest[e["provider"]]=e
    summary=" | ".join(f'{p}: HTTP {latest[p]["status"]} — {latest[p]["message"]}' for p in order if p in latest)
    raise HTTPException(503,{"message":"Қолжетімді AI мүмкіндігі аяқталды немесе провайдер жауап бермеді.","summary":summary,"errors":latest})

class ZerekChatBody(BaseModel):
    message:str
    provider:str="auto"
    mode:str="teacher"

@app.post("/api/zerek/chat")
async def zerek_chat(body:ZerekChatBody, request:Request):
    u=auth_user(request)
    role="Қазақстан мұғаліміне арналған қазақша AI көмекші" if body.mode=="teacher" else "Қазақстан оқушысына арналған түсіндіруші AI көмекші"
    prompt=f"""Сен — {role}. Пайдаланушыға нақты, пайдалы, құрылымды жауап бер.
Егер сұраныс ҚМЖ туралы болса, оқу мақсаты, сабақ мақсаты, кезеңдер, педагог әрекеті, оқушы әрекеті, бағалау, ресурстар логикасын ұсын.
Егер тапсырма/тест/дескриптор сұралса, бірден қолдануға дайын нұсқа бер.
Математикалық формулаларды LaTeX-пен жазуға болады. Жауап тілі пайдаланушы тіліне сай болсын.
Сұраныс: {body.message}"""
    ai=await ai_json(body.provider,prompt,u,"zerek_chat")
    data=ai["data"]
    if isinstance(data,dict):
        text=data.get("answer") or data.get("text") or data.get("response") or data.get("content") or json.dumps(data,ensure_ascii=False,indent=2)
    else:text=str(data)
    return {"ok":True,"text":text,"provider":ai["provider_used"],"fallback":ai["fallback"]}

@app.get("/api/ai/diagnostics")
def ai_diagnostics():
    return {"gemini":{"configured":bool(os.getenv("GEMINI_API_KEY","").strip()),"model":os.getenv("GEMINI_MODEL","gemini-3.6-flash")},"openai":{"configured":bool(os.getenv("OPENAI_API_KEY","").strip()),"model":os.getenv("OPENAI_MODEL","gpt-5.6")},"fallback":"selected -> retry x3 -> other provider -> retry x3"}

@app.get("/api/ai/status")
def ai_status():
    return {
        "gemini": bool(os.getenv("GEMINI_API_KEY","").strip()),
        "openai": bool(os.getenv("OPENAI_API_KEY","").strip()),
        "gemini_model": os.getenv("GEMINI_MODEL","gemini-2.5-flash"),
        "openai_model": os.getenv("OPENAI_MODEL","gpt-5.6")
    }


def _pick(d, *keys, default=""):
    if not isinstance(d, dict): return default
    for k in keys:
        v=d.get(k)
        if v not in (None,"",[],{}): return v
    return default

def normalize_task(t, i):
    if not isinstance(t, dict):
        t={"question":str(t)}
    options=_pick(t,"options","choices","variants",default=[])
    if not isinstance(options,list): options=[]
    return {
        "id": _pick(t,"id","number","task_number",default=i),
        "title": str(_pick(t,"title","name",default=f"{i}-тапсырма")),
        "question": str(_pick(t,"question","task","text","prompt",default="")),
        "answer": str(_pick(t,"answer","correct_answer","correctAnswer",default="")),
        "correct_answer": str(_pick(t,"correct_answer","correctAnswer","answer",default="")),
        "descriptor": str(_pick(t,"descriptor","descriptors","criterion","criteria",default="")),
        "type": str(_pick(t,"type","format",default="short")),
        "options": options,
        "max_score": _pick(t,"max_score","maxScore","score",default=1),
        "work_mode": str(_pick(t,"work_mode","workMode","mode",default="individual")),
        "group_name": str(_pick(t,"group_name","groupName","group",default="")),
        "external_url": str(_pick(t,"external_url","externalUrl","url","link",default=""))
    }

def normalize_lesson(raw):
    # Accept wrappers commonly returned by models.
    if isinstance(raw, list):
        raw={"tasks":raw}
    if not isinstance(raw, dict):
        raw={}
    for wrapper in ("data","lesson","result","analysis"):
        if isinstance(raw.get(wrapper),dict):
            merged=dict(raw)
            inner=merged.pop(wrapper)
            raw={**merged,**inner}
            break
    tasks=_pick(raw,"tasks","assignments","exercises","questions",default=[])
    if isinstance(tasks,dict): tasks=list(tasks.values())
    if not isinstance(tasks,list): tasks=[]
    return {
        "subject": str(_pick(raw,"subject","subject_name"," пән","пән",default="")),
        "class_name": str(_pick(raw,"class_name","class","grade","сынып",default="")),
        "topic": str(_pick(raw,"topic","lesson_topic","title","тақырып",default="")),
        "learning_objective": str(_pick(raw,"learning_objective","learningObjective","learning_objectives","оқу мақсаты",default="")),
        "lesson_goal": str(_pick(raw,"lesson_goal","lessonGoal","lesson_objective","сабақ мақсаты",default="")),
        "tasks": [normalize_task(t,i+1) for i,t in enumerate(tasks)]
    }

def normalize_generated(raw):
    if isinstance(raw,list): raw={"tasks":raw}
    if not isinstance(raw,dict): raw={}
    for wrapper in ("data","result","lesson"):
        if isinstance(raw.get(wrapper),dict):
            raw={**raw,**raw[wrapper]}
            break
    tasks=_pick(raw,"tasks","assignments","exercises","questions",default=[])
    if isinstance(tasks,dict): tasks=list(tasks.values())
    if not isinstance(tasks,list): tasks=[]
    return {"tasks":[normalize_task(t,i+1) for i,t in enumerate(tasks)]}

class KmjAnalyze(BaseModel):
    text: str
    hint: str = ""
    provider: str = "gemini"

@app.post("/api/kmj/analyze")
async def kmj_analyze(body: KmjAnalyze, request:Request):
    prompt=f"""Сен Қазақстан мұғаліміне арналған AI көмекшісің. ҚМЖ мәтінін мұқият талда.
Тек жарамды JSON қайтар: subject, class_name, topic, learning_objective, lesson_goal,
tasks массиві: id, title, question, answer, descriptor, type.
ҚМЖ-дағы нақты тапсырмаларды жоғалтпа. Жоқ жауапты бос қалдыр. Ойдан оқу мақсатын қоспа.
Мұғалім ескертпесі: {body.hint}
ҚМЖ:
{body.text[:50000]}"""
    try:
        ai=await ai_json(body.provider,prompt,auth_user(request),'kmj_analyze')
        raw=ai["data"]
        data=normalize_lesson(raw)
        if not any([data["subject"],data["class_name"],data["topic"],data["learning_objective"],data["tasks"]]):
            raise HTTPException(422,"AI жауап берді, бірақ ҚМЖ құрылымы анықталмады. ҚМЖ мәтіні/файлын тексеріңіз.")
        return {"ok":True,"mode":ai["provider_used"],"provider":ai["provider_used"],"fallback":ai["fallback"],"requested_provider":ai["requested_provider"],"attempts":ai["attempts"],"data":data}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500,f"AI талдау қатесі: {e}")

class TaskGenerate(BaseModel):
    lesson: dict
    format: str
    instruction: str = ""
    provider: str = "gemini"

@app.post("/api/tasks/generate")
async def task_generate(body: TaskGenerate, request:Request):
    prompt=f"""Сен мұғалім берген ҚМЖ негізінде интерактивті HTML тапсырма құрылымын дайындайсың.
Формат: {body.format}
Мұғалімнің міндетті ұсынысы: {body.instruction}
Тек жарамды JSON қайтар: {{"tasks":[...]}}.
Әр тапсырма: id,type,question,options(array),correct_answer,descriptor,max_score.
type тек test, match, fill, truefalse, short.
Оқу мақсатын сақта. ҚМЖ-да жоқ мазмұнды орынсыз қоспа.
ҚМЖ дерегі:
{json.dumps(body.lesson,ensure_ascii=False)[:40000]}"""
    try:
        ai=await ai_json(body.provider,prompt,auth_user(request),'task_generate')
        raw=ai["data"]
        data=normalize_generated(raw)
        if not data["tasks"]:
            raise HTTPException(422,"AI тапсырма құрылымын қайтармады. ҚМЖ-дан тапсырма табылғанын және нұсқауды тексеріңіз.")
        return {"ok":True,"mode":ai["provider_used"],"provider":ai["provider_used"],"fallback":ai["fallback"],"requested_provider":ai["requested_provider"],"attempts":ai["attempts"],"data":data}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500,f"AI тапсырма жасау қатесі: {e}")

class StudentAnalyze(BaseModel):
    provider: str = "gemini"
    lesson: dict
    student: dict

@app.post("/api/student/analyze")
async def student_analyze(body: StudentAnalyze, request:Request):
    prompt=f"""Сен мұғалімге арналған көмекші AI-сің. Оқушыға баға қойма.
Оқу мақсатына жетуін дәлелдермен талда және мұғалімге қысқа педагогикалық қорытынды бер.
Камерадағы қозғалыс көрсеткішін зейін, эмоция немесе оқу жетістігі деп түсіндірме.
Тек JSON қайтар: summary, objective_status, strengths(array), needs_support(array), evidence(array), next_step.
Сабақ: {json.dumps(body.lesson,ensure_ascii=False)}
Оқушы дерегі: {json.dumps(body.student,ensure_ascii=False)}
"""
    ai=await ai_json(body.provider,prompt,auth_user(request),'student_analyze')
    return {"mode":ai["provider_used"],"fallback":ai["fallback"],"data":ai["data"]}

@app.post("/api/audio/transcribe")
async def audio_transcribe(file: UploadFile=File(...), provider: str="openai"):
    data=await file.read()
    provider=(provider or "openai").lower()
    if provider=="openai":
        key=os.getenv("OPENAI_API_KEY","").strip()
        if not key: raise HTTPException(400,"OPENAI_API_KEY орнатылмаған.")
        import httpx
        model=os.getenv("OPENAI_TRANSCRIBE_MODEL","gpt-4o-mini-transcribe")
        files={"file":(file.filename or "answer.webm",data,file.content_type or "audio/webm")}
        form={"model":model,"language":"kk"}
        async with httpx.AsyncClient(timeout=120) as client:
            r=await client.post("https://api.openai.com/v1/audio/transcriptions",headers={"Authorization":f"Bearer {key}"},data=form,files=files)
            if r.status_code>=400: raise HTTPException(r.status_code,f"OpenAI transcription: {r.text[:500]}")
            return {"provider":"openai","text":r.json().get("text","")}
    key=os.getenv("GEMINI_API_KEY","").strip()
    if not key: raise HTTPException(400,"GEMINI_API_KEY орнатылмаған.")
    import httpx
    model=os.getenv("GEMINI_MODEL","gemini-2.5-flash")
    url=f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={key}"
    mime=file.content_type or "audio/webm"
    payload={"contents":[{"parts":[
        {"text":"Осы аудиодағы қазақша сөйлеуді дәл мәтінге түсір. Тек транскрипт мәтінін қайтар."},
        {"inline_data":{"mime_type":mime,"data":base64.b64encode(data).decode("ascii")}}
    ]}]}
    async with httpx.AsyncClient(timeout=120) as client:
        r=await client.post(url,json=payload)
        if r.status_code>=400: raise HTTPException(r.status_code,f"Gemini audio: {r.text[:500]}")
        text=r.json()["candidates"][0]["content"]["parts"][0]["text"]
        return {"provider":"gemini","text":text.strip()}


URL_RE = re.compile(r'https?://[^\s<>"\')\]]+')

class ResourceExtract(BaseModel):
    text: str

@app.post("/api/resources/extract")
def extract_resources(body: ResourceExtract):
    urls=[]
    for u in URL_RE.findall(body.text or ""):
        u=u.rstrip(".,;:")
        if u not in urls: urls.append(u)
    return {"resources":[{"url":u,"approved":False} for u in urls[:30]]}


class DesignAnalyze(BaseModel):
    provider: str = "gemini"
    image_data_url: str
    instruction: str = ""

@app.post("/api/design/analyze")
async def design_analyze(body: DesignAnalyze):
    import base64 as _b64, httpx
    if "," not in body.image_data_url:
        raise HTTPException(400,"Үлгі сурет оқылмады.")
    meta,b64=body.image_data_url.split(",",1)
    mime=(meta.split(";")[0].split(":")[-1] or "image/png")
    instruction=body.instruction or "Үлгінің жалпы орналасуын, түстерін және карточка құрылымын сипатта."
    prompt="""Бұл мұғалім жіберген жұмыс парағының дизайн үлгісі.
Суреттегі тапсырма мәтінін көшірме. Тек визуалдық құрылымды талда.
Тек JSON қайтар:
{"layout":"grid2x2|grid1x4|grid2x3|free","columns":2,"card_count":4,
"palette":["#hex"],"background":"#hex","title_style":"pill|plain|banner",
"border_radius":18,"border_width":3,"font_scale":1.0,
"image_position":"right|left|top|none","notes":"қысқа сипаттама"}.
Мұғалім нұсқауы: """+instruction
    provider=(body.provider or "gemini").lower()
    if provider=="openai":
        key=os.getenv("OPENAI_API_KEY","").strip()
        if not key: raise HTTPException(400,"OPENAI_API_KEY жоқ.")
        model=os.getenv("OPENAI_MODEL","gpt-5.6")
        payload={"model":model,"input":[{"role":"user","content":[
            {"type":"input_text","text":prompt},
            {"type":"input_image","image_url":body.image_data_url}
        ]}]}
        async with httpx.AsyncClient(timeout=120) as client:
            r=await client.post("https://api.openai.com/v1/responses",headers={"Authorization":f"Bearer {key}","Content-Type":"application/json"},json=payload)
            if r.status_code>=400: raise HTTPException(r.status_code,f"OpenAI vision: {r.text[:500]}")
            d=r.json(); text=""
            for item in d.get("output",[]):
                if item.get("type")=="message":
                    for c in item.get("content",[]):
                        if c.get("type")=="output_text": text+=c.get("text","")
            return {"provider":"openai","design":_json_from_text(text)}
    key=os.getenv("GEMINI_API_KEY","").strip()
    if not key: raise HTTPException(400,"GEMINI_API_KEY жоқ.")
    model=os.getenv("GEMINI_MODEL","gemini-3.6-flash")
    url=f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={key}"
    payload={"contents":[{"parts":[{"text":prompt},{"inline_data":{"mime_type":mime,"data":b64}}]}],"generationConfig":{"responseMimeType":"application/json"}}
    async with httpx.AsyncClient(timeout=120) as client:
        r=await client.post(url,json=payload)
        if r.status_code>=400: raise HTTPException(r.status_code,f"Gemini vision: {r.text[:500]}")
        text=r.json()["candidates"][0]["content"]["parts"][0]["text"]
        return {"provider":"gemini","design":_json_from_text(text)}

# ---------------- Lesson / task sessions ----------------
LESSONS={}
class PublishLesson(BaseModel):
    subject:str=""
    class_name:str=""
    topic:str=""
    learning_objective:str=""
    tasks:list=[]

@app.post("/api/lesson/publish")
def publish_lesson(body: PublishLesson):
    code=secrets.token_hex(4).upper()
    LESSONS[code]={"subject":body.subject,"class_name":body.class_name,"topic":body.topic,"learning_objective":body.learning_objective,"tasks":body.tasks,"results":[]}
    return {"code":code}

@app.get("/api/lesson/{code}")
def get_lesson(code:str):
    if code not in LESSONS:raise HTTPException(404,"Сабақ коды табылмады.")
    return LESSONS[code]

class TaskSubmit(BaseModel):
    code:str
    student_id:str=""
    name:str
    answers:dict

@app.post("/api/lesson/submit")
def submit_tasks(body:TaskSubmit):
    if body.code not in LESSONS:raise HTTPException(404,"Сабақ табылмады.")
    lesson=LESSONS[body.code]; items=[]; score=0; total=0
    for t in lesson["tasks"]:
        tid=str(t.get("id","")); ans=str(body.answers.get(tid,"")).strip()
        correct=str(t.get("correct_answer","")).strip()
        maxs=float(t.get("max_score",1) or 1); total+=maxs
        ok=bool(correct) and ans.lower()==correct.lower()
        got=maxs if ok else 0;score+=got
        items.append({"id":tid,"answer":ans,"correct_answer":correct,"correct":ok,"score":got,"max_score":maxs,"descriptor":t.get("descriptor","")})
    result={"student_id":body.student_id,"name":body.name,"score":score,"total":total,"items":items,"time":datetime.now().isoformat()}
    lesson["results"].append(result)
    return result

# ---------------- QR attendance ----------------
import qrcode
QR_SESSIONS={}
class SessionStart(BaseModel):
    class_name:str="Сынып"; topic:str="Сабақ"; minutes:int=10
class CheckIn(BaseModel):
    token:str; student_id:str=""; name:str; class_name:str

@app.post("/api/session/start")
def session_start(body:SessionStart):
    token=secrets.token_hex(5).upper();now=datetime.now();minutes=max(1,min(body.minutes,120))
    QR_SESSIONS[token]={"class_name":body.class_name,"topic":body.topic,"active":True,"expires_at":now+timedelta(minutes=minutes),"attendance":{}}
    s=QR_SESSIONS[token];return {"token":token,"class_name":s["class_name"],"topic":s["topic"],"expires_at":s["expires_at"].isoformat()}
def qs(token):
    s=QR_SESSIONS.get(token)
    if not s:raise HTTPException(404,"QR-сабақ табылмады.")
    if datetime.now()>=s["expires_at"]:s["active"]=False
    return s
@app.get("/api/session/{token}")
def session_info(token:str):
    s=qs(token);return {"class_name":s["class_name"],"topic":s["topic"],"active":s["active"],"expires_at":s["expires_at"].isoformat()}
@app.get("/api/session/{token}/qr.png")
def session_qr(token:str,base:str):
    qs(token);img=qrcode.make(base.rstrip("/")+"/?session="+token);buf=BytesIO();img.save(buf,format="PNG");return Response(buf.getvalue(),media_type="image/png")
@app.post("/api/session/stop/{token}")
def session_stop(token:str):
    s=qs(token);s["active"]=False;return {"ok":True}
@app.post("/api/checkin")
def checkin(body:CheckIn):
    s=qs(body.token)
    if not s["active"]:raise HTTPException(410,"QR аяқталған.")
    key=(body.student_id.strip() or (body.name+"|"+body.class_name)).lower()
    if key not in s["attendance"]:s["attendance"][key]={"student_id":body.student_id.strip().upper(),"name":body.name.strip(),"class_name":body.class_name.strip(),"time":datetime.now().strftime("%H:%M:%S")}
    return {"ok":True,"time":s["attendance"][key]["time"]}
@app.get("/api/session/{token}/attendance")
def session_attendance(token:str):
    s=qs(token);return {"class_name":s["class_name"],"topic":s["topic"],"active":s["active"],"attendance":list(s["attendance"].values())}

@app.get("/favicon.png", include_in_schema=False)
def favicon(): return FileResponse(BASE/"favicon.png")

@app.get("/zerek-logo.png", include_in_schema=False)
def zerek_logo(): return FileResponse(BASE/"zerek-logo.png")

@app.get("/manifest.webmanifest",include_in_schema=False)
def pwa_manifest():return FileResponse(BASE/"manifest.webmanifest",media_type="application/manifest+json")
@app.get("/sw.js",include_in_schema=False)
def sw():return FileResponse(BASE/"sw.js",media_type="application/javascript")
@app.get("/robots.txt",include_in_schema=False)
def robots():return FileResponse(BASE/"robots.txt",media_type="text/plain")
@app.get("/sitemap.xml",include_in_schema=False)
def sitemap():return FileResponse(BASE/"sitemap.xml",media_type="application/xml")
@app.get("/icon-180.png",include_in_schema=False)
def i180():return FileResponse(BASE/"icon-180.png")
@app.get("/icon-192.png",include_in_schema=False)
def i192():return FileResponse(BASE/"icon-192.png")
@app.get("/icon-512.png",include_in_schema=False)
def i512():return FileResponse(BASE/"icon-512.png")
@app.get("/")
def home():return FileResponse(BASE/"index.html")
app.mount("/static",StaticFiles(directory=BASE),name="static")
